"""Extract text from downloaded UFLI resources (PPTX and PDF)."""

from __future__ import annotations

import json
import logging
from dataclasses import asdict, dataclass
from pathlib import Path

logger = logging.getLogger(__name__)


@dataclass
class LessonContent:
    """Extracted content for one UFLI lesson."""

    lesson_id: str
    lesson_group: str
    concept: str
    slide_text: str
    slide_count: int
    decodable_text: str
    home_practice_text: str
    additional_text: str


def extract_pptx_text(path: str) -> tuple[str, int]:
    """Extract all text from a PPTX file.

    Returns:
        Tuple of (combined text, slide count).
    """
    from pptx import Presentation

    prs = Presentation(path)
    texts: list[str] = []
    for slide in prs.slides:
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_parts.append(shape.text_frame.text)
        if slide_parts:
            texts.append("\n".join(slide_parts))
    return "\n\n".join(texts), len(prs.slides)


def extract_pdf_text(path: str) -> str:
    """Extract text from a PDF file using PyMuPDF."""
    import fitz

    doc = fitz.open(path)
    pages = [page.get_text() for page in doc]
    doc.close()
    return "\n\n".join(pages)


def extract_lesson(
    lesson_id: str,
    lesson_group: str,
    concept: str,
    data_dir: str = "data/ufli",
) -> LessonContent | None:
    """Extract content for one lesson from downloaded files."""
    lesson_dir = Path(data_dir) / "raw" / lesson_id
    if not lesson_dir.exists():
        return None

    slide_text = ""
    slide_count = 0
    decodable_text = ""
    home_practice_text = ""
    additional_text = ""

    # PPTX slide deck (prefer direct download, fall back to Google Slides export)
    pptx_path = lesson_dir / "slide_deck_pptx.pptx"
    if not pptx_path.exists():
        pptx_path = lesson_dir / "slide_deck_gslides_export.pptx"
    if pptx_path.exists():
        try:
            slide_text, slide_count = extract_pptx_text(str(pptx_path))
        except Exception:
            logger.exception("Failed to extract PPTX for lesson %s", lesson_id)

    # Decodable passage PDF
    decodable_path = lesson_dir / "decodable_passage_pdf.pdf"
    if decodable_path.exists():
        try:
            decodable_text = extract_pdf_text(str(decodable_path))
        except Exception:
            logger.exception("Failed to extract decodable PDF for lesson %s", lesson_id)

    # Home practice PDF
    home_path = lesson_dir / "home_practice_pdf.pdf"
    if home_path.exists():
        try:
            home_practice_text = extract_pdf_text(str(home_path))
        except Exception:
            logger.exception("Failed to extract home practice PDF for lesson %s", lesson_id)

    # Additional activities PDF
    additional_path = lesson_dir / "additional_pdf.pdf"
    if additional_path.exists():
        try:
            additional_text = extract_pdf_text(str(additional_path))
        except Exception:
            logger.exception("Failed to extract additional PDF for lesson %s", lesson_id)

    if not slide_text and not decodable_text and not home_practice_text:
        logger.warning("No extractable content for lesson %s", lesson_id)
        return None

    return LessonContent(
        lesson_id=lesson_id,
        lesson_group=lesson_group,
        concept=concept,
        slide_text=slide_text,
        slide_count=slide_count,
        decodable_text=decodable_text,
        home_practice_text=home_practice_text,
        additional_text=additional_text,
    )


def extract_all(data_dir: str = "data/ufli") -> list[LessonContent]:
    """Extract all downloaded lessons and write normalized.jsonl."""
    base = Path(data_dir)
    manifest_path = base / "manifest.jsonl"
    if not manifest_path.exists():
        logger.error("No manifest.jsonl found in %s", data_dir)
        return []

    results: list[LessonContent] = []
    for line in manifest_path.read_text().splitlines():
        if not line.strip():
            continue
        rec = json.loads(line)
        content = extract_lesson(
            lesson_id=rec["lesson_id"],
            lesson_group=rec.get("lesson_group", ""),
            concept=rec.get("concept", ""),
            data_dir=data_dir,
        )
        if content:
            results.append(content)

    # Write normalized.jsonl
    normalized_path = base / "normalized.jsonl"
    with normalized_path.open("w") as f:
        for item in results:
            f.write(json.dumps(asdict(item)) + "\n")

    logger.info("Extracted %d lessons to %s", len(results), normalized_path)
    return results
