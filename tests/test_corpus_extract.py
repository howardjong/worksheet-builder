"""Tests for UFLI corpus text extraction."""

from __future__ import annotations

from pathlib import Path

import pytest

from corpus.ufli.extract import extract_pptx_text


def test_extract_pptx_text(tmp_path: Path) -> None:
    """PPTX text extraction returns slide text and count."""
    pptx = pytest.importorskip("pptx")

    prs = pptx.Presentation()
    layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(layout)
    from pptx.util import Inches

    tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tx_box.text_frame.text = "CVCe Words: grade, chase, slide"

    path = tmp_path / "test.pptx"
    prs.save(str(path))

    text, count = extract_pptx_text(str(path))
    assert "grade" in text
    assert "chase" in text
    assert count == 1


def test_extract_pptx_multiple_slides(tmp_path: Path) -> None:
    """Multiple slides produce combined text with correct count."""
    pptx = pytest.importorskip("pptx")

    prs = pptx.Presentation()
    layout = prs.slide_layouts[6]
    from pptx.util import Inches

    for word in ["apple", "banana"]:
        slide = prs.slides.add_slide(layout)
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tx_box.text_frame.text = word

    path = tmp_path / "multi.pptx"
    prs.save(str(path))

    text, count = extract_pptx_text(str(path))
    assert "apple" in text
    assert "banana" in text
    assert count == 2


def test_extract_pdf_text(tmp_path: Path) -> None:
    """PDF text extraction returns page content."""
    fitz = pytest.importorskip("fitz")
    from corpus.ufli.extract import extract_pdf_text

    # Create a minimal PDF with PyMuPDF
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Decodable passage: The cat sat on the mat.")
    path = tmp_path / "test.pdf"
    doc.save(str(path))
    doc.close()

    text = extract_pdf_text(str(path))
    assert "cat sat" in text
